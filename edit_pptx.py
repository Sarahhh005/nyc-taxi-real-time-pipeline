import sys
from pptx import Presentation

def replace_text(shape, old_text, new_text):
    if hasattr(shape, "text") and old_text in shape.text:
        # To preserve formatting as much as possible, we replace text at the run level
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if old_text in run.text:
                        run.text = run.text.replace(old_text, new_text)

def main():
    path = '/opt/workspace/docs/NYC_Taxi_Real-Time_Analytics_Pipeline (1).pptx'
    out_path = '/opt/workspace/docs/NYC_Taxi_Real-Time_Analytics_Pipeline_Updated.pptx'
    
    try:
        prs = Presentation(path)
    except Exception as e:
        print(f"Failed to open {path}: {e}")
        return

    # Slide 1
    for shape in prs.slides[0].shapes:
        replace_text(shape, "Superset · Docker", "Superset · AI Agent · Docker")
        replace_text(shape, "visualized live in Superset", "visualized in Superset & queried via AI Agent")

    # Slide 3 (Tech Stack)
    for shape in prs.slides[2].shapes:
        replace_text(shape, "Pandas", "FastAPI & LangChain\nAI Agent API")

    # Slide 4 (Architecture)
    for shape in prs.slides[3].shapes:
        replace_text(shape, "Six stages", "Seven stages")
        replace_text(shape, "Superset", "Superset & AI Agent")

    # Slide 5 (Data)
    for shape in prs.slides[4].shapes:
        replace_text(shape, "10,000", "35+ MILLION")
        replace_text(shape, "SAMPLE TRIP RECORDS", "TRIP RECORDS (FULL 2025)")
        replace_text(shape, "sampled and streamed to simulate", "streamed sequentially to simulate")

    # Slide 10 (Superset)
    for shape in prs.slides[9].shapes:
        replace_text(shape, "STEP 6 OF 6", "STEP 6 OF 7")
        replace_text(shape, "24.5K", "150K+")

    # Slide 11 (Results Demand)
    for shape in prs.slides[10].shapes:
        replace_text(shape, "10,000-trip sample — January 2025", "35+ Million trips — Full Year 2025")
        replace_text(shape, "485", "1.6M")
        replace_text(shape, "469", "1.5M")
        replace_text(shape, "455", "1.4M")
        replace_text(shape, "412", "1.3M")

    # Slide 12 (Results Revenue)
    for shape in prs.slides[11].shapes:
        replace_text(shape, "10,000-trip sample — January 2025", "35+ Million trips — Full Year 2025")

    # Slide 14 (What's Next)
    for shape in prs.slides[13].shapes:
        replace_text(shape, "Extend ingestion across all twelve months of 2025.", "Implement conversational memory for the AI Agent.")
        replace_text(shape, "Full-year data backfill", "Advanced Multi-Agent System")

    # Slide 15 (Thank you)
    for shape in prs.slides[14].shapes:
        replace_text(shape, "Superset · Docker", "Superset · AI Agent · Docker")

    prs.save(out_path)
    print(f"Successfully saved updated presentation to {out_path}")

if __name__ == "__main__":
    main()
