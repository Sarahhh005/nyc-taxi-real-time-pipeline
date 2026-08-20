from pptx import Presentation
prs1 = Presentation('/opt/workspace/docs/NYC_Taxi_Real-Time_Analytics_Pipeline (1).pptx')
prs2 = Presentation('/opt/workspace/docs/NYC_Taxi_Real-Time_Analytics_Pipeline.pptx')
print('Slide count 1:', len(prs1.slides))
print('Slide count 2:', len(prs2.slides))

print('\n--- Slides in prs2 ---')
for i, slide in enumerate(prs2.slides):
    print(f"\n--- Slide {i+1} ---")
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            print(shape.text)
