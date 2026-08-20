import sys
from pptx import Presentation

def main():
    path = '/opt/workspace/docs/NYC_Taxi_Real-Time_Analytics_Pipeline (1).pptx'
    try:
        prs = Presentation(path)
    except Exception as e:
        print(f"Failed to open {path}: {e}")
        return

    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                print(shape.text)

if __name__ == "__main__":
    main()
